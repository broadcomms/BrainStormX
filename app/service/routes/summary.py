# app/service/routes/summary.py
from __future__ import annotations

import json
import os
from datetime import datetime
from typing import Any, Dict, List, Optional, Tuple

from flask import current_app
from sqlalchemy import func
from sqlalchemy.orm import joinedload

from app.extensions import db
from app.config import Config
from app.models import (
    Workshop,
    BrainstormTask,
    BrainstormIdea,
    IdeaCluster,
    IdeaVote,
    ChatMessage,
    Transcript,
    Document,
    WorkshopDocument,
)
from app.utils.json_utils import extract_json_block
from app.utils.data_aggregation import get_pre_workshop_context_json
from app.utils.llm_bedrock import get_chat_llm_pro
from langchain_core.prompts import PromptTemplate


# =============== Utilities ===============

class SummaryGenerationError(RuntimeError):
    pass

def _safe(s: Any) -> str:
    return (str(s) if s is not None else "").strip()

def _coerce_text(val: Any) -> str:
    if val is None:
        return ""
    if isinstance(val, str):
        return val
    if isinstance(val, dict):
        if "content" in val and isinstance(val["content"], str):
            return val["content"]
        return json.dumps(val, ensure_ascii=False)
    if hasattr(val, "content"):
        c = getattr(val, "content", None)
        if isinstance(c, str):
            return c
    return str(val)

def _reports_dir() -> str:
    base = os.path.join(current_app.instance_path, "uploads", "reports")
    os.makedirs(base, exist_ok=True)
    return base

def _media_url(rel_path: str) -> str:
    return f"{Config.MEDIA_REPORTS_URL_PREFIX}/{os.path.basename(rel_path)}"

def _safe_title(s: str) -> str:
    return (s or "Workshop").strip().replace("/", "-")


# =============== Artifact Generators (LLM-owned specs only; no content fallbacks) ===============

def _generate_pdf_from_doc_spec(ws: Workshop, doc_spec: Dict[str, Any]) -> Tuple[str, str, str] | None:
    """Render a professional PDF from an LLM-provided document_spec.
    The spec schema matches feasibility’s JSON renderer and supports:
      - title (str)
      - cover: {subtitle, objective, date_str, top_clusters:[{name,votes}]}
      - sections: [{heading, blocks:[{type: p|h2|ul|table|note|rule|page_break, ...}]}]
      - appendices: same shape as sections (optional)
    """
    try:
        from reportlab.lib import colors
        from reportlab.lib.enums import TA_LEFT
        from reportlab.lib.pagesizes import LETTER
        from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
        from reportlab.lib.units import inch
        from reportlab.platypus import (
            BaseDocTemplate, Frame, PageTemplate, Paragraph, Spacer, Table, TableStyle,
            ListFlowable, ListItem, PageBreak
        )
        from reportlab.pdfgen import canvas
    except Exception as exc:
        current_app.logger.warning("[Summary] ReportLab unavailable: %s", exc)
        return None

    ts = datetime.utcnow().strftime("%Y-%m-%d %H-%M-%S")
    fname = f"{_safe_title(ws.title)} summary {ts}.pdf"
    abs_path = os.path.join(_reports_dir(), fname)
    rel_path = os.path.join("uploads", "reports", fname)

    base = getSampleStyleSheet()
    Title = ParagraphStyle("BX_Title", parent=base["Title"], fontSize=24, leading=28, alignment=TA_LEFT,
                           spaceBefore=18, spaceAfter=10, textColor=colors.HexColor("#111827"))
    Sub = ParagraphStyle("BX_Sub", parent=base["Heading2"], fontSize=12, leading=16,
                         textColor=colors.HexColor("#4B5563"), spaceAfter=10)
    H1 = ParagraphStyle("BX_H1", parent=base["Heading2"], fontSize=16, leading=20,
                        spaceBefore=16, spaceAfter=8, textColor=colors.HexColor("#111827"))
    H2 = ParagraphStyle("BX_H2", parent=base["Heading3"], fontSize=13, leading=17,
                        spaceBefore=10, spaceAfter=6, textColor=colors.HexColor("#111827"))
    Body = ParagraphStyle("BX_Body", parent=base["BodyText"], fontSize=10.5, leading=14.0,
                          textColor=colors.HexColor("#111827"), spaceAfter=6)
    Note = ParagraphStyle("BX_Note", parent=base["BodyText"], fontSize=9.5, leading=13,
                          textColor=colors.HexColor("#6B7280"), spaceBefore=4, spaceAfter=6)

    def rule(width=6.5 * inch, height=0.7, color="#E5E7EB", space=8):
        t = Table([[""]], colWidths=[width], rowHeights=[height])
        t.setStyle(TableStyle([("BACKGROUND", (0, 0), (-1, -1), colors.HexColor(color))]))
        return [Spacer(1, space), t, Spacer(1, space)]

    def _footer(canv: canvas.Canvas, doc):
        canv.setFont("Helvetica", 9)
        canv.setFillColor(colors.HexColor("#6B7280"))
        canv.drawRightString(7.95 * inch, 0.5 * inch, f"Page {doc.page}")

    frame = Frame(0.75 * inch, 0.75 * inch, 7.0 * inch, 9.75 * inch, showBoundary=0)
    doc = BaseDocTemplate(abs_path, pagesize=LETTER, leftMargin=0, rightMargin=0, topMargin=0, bottomMargin=0)
    doc.addPageTemplates([PageTemplate(id="main", frames=[frame], onPage=_footer)])

    def _ul(items: List[str]) -> ListFlowable:
        li = [ListItem(Paragraph(_safe(i), Body)) for i in (items or []) if _safe(i)]
        return ListFlowable(li, bulletType="bullet", leftIndent=12)

    def _table(columns: List[str], rows: List[List[str]]) -> Table:
        cols = [Paragraph(_safe(c), Body) for c in (columns or [])]
        data = [cols] + [[Paragraph(_safe(c), Body) for c in (r or [])] for r in (rows or [])]
        tbl = Table(data, colWidths=[(6.5 * inch) / max(1, len(cols))] * max(1, len(cols)))
        tbl.setStyle(
            TableStyle(
                [
                    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#F3F4F6")),
                    ("GRID", (0, 0), (-1, -1), 0.25, colors.HexColor("#D1D5DB")),
                    ("ALIGN", (0, 0), (-1, -1), "LEFT"),
                    ("VALIGN", (0, 0), (-1, -1), "TOP"),
                    ("TOPPADDING", (0, 0), (-1, -1), 4),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
                ]
            )
        )
        return tbl

    E: List[Any] = []

    # Cover
    title_text = _safe(doc_spec.get("title") or f"{ws.title} — Executive Summary")
    E.append(Paragraph(title_text, Title))
    cover = doc_spec.get("cover") or {}
    subtitle = _safe(cover.get("subtitle"))
    if subtitle:
        E.append(Paragraph(subtitle, Sub))
    E.extend(rule(color="#0d6efd"))
    objective = _safe(cover.get("objective") or getattr(ws, "objective", ""))
    if objective:
        E.append(Paragraph("Objective", H2))
        E.append(Paragraph(objective, Body))
    date_str = _safe(cover.get("date_str") or (ws.date_time.strftime("%Y-%m-%d %H:%M UTC") if getattr(ws, "date_time", None) else "TBD"))
    E.append(Paragraph("Date", H2))
    E.append(Paragraph(date_str, Body))
    top_clusters = cover.get("top_clusters") or []
    if top_clusters:
        E.append(Paragraph("Top Voted Clusters", H2))
        cols = ["Cluster", "Votes"]
        rows = [[_safe(c.get("name")), str(int(c.get("votes") or 0))] for c in top_clusters]
        E.append(_table(cols, rows))
    E.append(PageBreak())

    # Sections & Appendices
    def _emit_sections(sections: List[Dict[str, Any]]):
        for sec in sections or []:
            heading = _safe(sec.get("heading"))
            if heading:
                E.append(Paragraph(heading, H1))
            for blk in sec.get("blocks", []):
                btype = (blk.get("type") or "p").strip().lower()
                if btype == "p":
                    E.append(Paragraph(_safe(blk.get("text")), Body))
                elif btype == "h2":
                    E.append(Paragraph(_safe(blk.get("text")), H2))
                elif btype == "ul":
                    E.append(_ul([_safe(x) for x in blk.get("items", [])]))
                elif btype == "table":
                    E.append(_table(blk.get("columns") or [], blk.get("rows") or []))
                elif btype == "note":
                    E.append(Paragraph(_safe(blk.get("text")), Note))
                elif btype == "rule":
                    E.extend(rule())
                elif btype == "page_break":
                    E.append(PageBreak())

    _emit_sections(doc_spec.get("sections") or [])
    if doc_spec.get("appendices"):
        E.append(PageBreak())
        E.append(Paragraph("Appendices", H1))
        _emit_sections(doc_spec.get("appendices") or [])

    E.extend(rule())
    prepared = datetime.utcnow().strftime("%Y-%m-%d %H:%M UTC")
    E.append(Paragraph(f"Prepared by BrainStormX • {prepared}", Note))

    try:
        doc.build(E)
    except Exception as exc:
        current_app.logger.error("[Summary] PDF build failed: %s", exc, exc_info=True)
        return None

    url = _media_url(rel_path)
    return abs_path, rel_path, url


def _generate_pptx_from_slides_spec(ws: Workshop, slides_spec: Dict[str, Any]) -> Tuple[str, str, str] | None:
    """Build a simple PPTX deck from an LLM-provided slides_spec:
       {
         "title": "...",
         "slides": [
           {"layout":"title+bullets", "title":"...", "bullets":["...","..."]},
           {"layout":"title+image", "title":"...", "image_url": "..."},
           {"layout":"title+table", "title":"...", "columns":[...], "rows":[...]},
           ...
         ]
       }
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except Exception as exc:
        current_app.logger.warning("[Summary] python-pptx unavailable: %s", exc)
        return None

    ts = datetime.utcnow().strftime("%Y-%m-%d %H-%M-%S")
    fname = f"{_safe_title(ws.title)} summary {ts}.pptx"
    abs_path = os.path.join(_reports_dir(), fname)
    rel_path = os.path.join("uploads", "reports", fname)

    prs = Presentation()
    # Title slide
    title = _safe(slides_spec.get("title") or f"{ws.title} — Workshop Summary")
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    subtitle = slide.placeholders[1]
    subtitle.text = _safe(getattr(ws, "objective", "") or "Executive brief and decisions")

    for s in slides_spec.get("slides", []):
        layout = (s.get("layout") or "title+bullets").lower()
        if layout == "title+bullets":
            sl = prs.slides.add_slide(prs.slide_layouts[1])
            sl.shapes.title.text = _safe(s.get("title"))
            body = sl.shapes.placeholders[1].text_frame
            body.clear()
            for b in s.get("bullets", []):
                p = body.add_paragraph() if body.paragraphs else body.paragraphs[0]
                p.text = _safe(b)
                p.level = 0
        elif layout == "title+table":
            sl = prs.slides.add_slide(prs.slide_layouts[5])
            sl.shapes.title.text = _safe(s.get("title"))
            cols = s.get("columns") or []
            rows = s.get("rows") or []
            x, y, cx, cy = Inches(1), Inches(1.8), Inches(8), Inches(4.5)
            table = sl.shapes.add_table(len(rows) + 1, max(1, len(cols)), x, y, cx, cy).table
            for i, c in enumerate(cols):
                table.cell(0, i).text = _safe(c)
            for r_idx, r in enumerate(rows, start=1):
                for c_idx, val in enumerate(r[: len(cols)]):
                    table.cell(r_idx, c_idx).text = _safe(val)
        else:
            # fallback to bullets layout if unknown (structure only; still LLM-owned content)
            sl = prs.slides.add_slide(prs.slide_layouts[1])
            sl.shapes.title.text = _safe(s.get("title"))
            body = sl.shapes.placeholders[1].text_frame
            body.clear()
            for b in s.get("bullets", []):
                p = body.add_paragraph() if body.paragraphs else body.paragraphs[0]
                p.text = _safe(b)
                p.level = 0

    try:
        prs.save(abs_path)
    except Exception as exc:
        current_app.logger.error("[Summary] PPTX build failed: %s", exc, exc_info=True)
        return None
    url = _media_url(rel_path)
    return abs_path, rel_path, url


def _attach_doc(ws: Workshop, abs_path: str, rel_path: str, url: str, payload: Dict[str, Any], title_suffix: str) -> None:
    """Persist Document and link to workshop; mirror into payload."""
    try:
        try:
            size = os.path.getsize(abs_path)
        except OSError:
            size = None
        doc = Document()
        doc.workspace_id = ws.workspace_id
        doc.title = f"{ws.title} — {title_suffix}"
        doc.description = f"Automatically generated {title_suffix.lower()}"
        doc.file_name = os.path.basename(rel_path)
        doc.file_path = rel_path
        doc.uploaded_by_id = ws.created_by_id
        doc.file_size = size
        db.session.add(doc)
        db.session.flush()

        link = WorkshopDocument()
        link.workshop_id = ws.id
        link.document_id = doc.id
        db.session.add(link)
        db.session.flush()

        doc_payload = {
            "id": doc.id,
            "title": doc.title,
            "file_name": doc.file_name,
            "file_size": size,
            "file_path": rel_path,
            "url": url,
            "workshop_link_id": link.id,
        }
        payload.setdefault("documents", {})
        payload["documents"][title_suffix.lower().replace(" ", "_")] = dict(doc_payload)
    except Exception as exc:
        current_app.logger.warning("[Summary] Skipped document attachment: %s", exc)


# =============== Input assembly ===============

def _load_latest_payload(workshop_id: int, types: List[str]) -> Optional[Dict[str, Any]]:
    try:
        task = (
            BrainstormTask.query
            .filter(BrainstormTask.workshop_id == workshop_id, BrainstormTask.task_type.in_(types))
            .order_by(BrainstormTask.created_at.desc())
            .first()
        )
        if not task or not task.payload_json:
            return None
        data = json.loads(task.payload_json)
        return data if isinstance(data, dict) else None
    except Exception:
        current_app.logger.debug("[Summary] Failed to load payload types=%s", types, exc_info=True)
        return None

def _collect_full_context(workshop_id: int) -> Dict[str, Any]:
    ws = db.session.get(Workshop, workshop_id)
    if not ws:
        raise SummaryGenerationError(f"Workshop {workshop_id} not found")

    # Pre-work + basic snapshot
    pre = get_pre_workshop_context_json(workshop_id)
    overview = {
        "title": ws.title,
        "objective": ws.objective or "TBD",
        "scheduled_for": ws.date_time.strftime("%Y-%m-%d %H:%M UTC") if ws.date_time else "unscheduled",
        "duration_minutes": ws.duration,
        "status": ws.status,
        "organizer": getattr(ws.creator, "display_name", None) or getattr(ws.creator, "email", None) or "Unknown organizer",
        "participant_count": (ws.participants.count() if hasattr(ws.participants, "count") else len(list(ws.participants))) if hasattr(ws, "participants") else 0,
    }

    # Phase payloads
    framing = _load_latest_payload(workshop_id, ["framing"]) or {}
    warmup = _load_latest_payload(workshop_id, ["warm-up", "warm_up", "introduction"]) or {}
    brainstorming = _load_latest_payload(workshop_id, ["brainstorming", "ideas"]) or {}
    clustering_voting = _load_latest_payload(workshop_id, ["clustering_voting"]) or {}
    feasibility = _load_latest_payload(workshop_id, ["results_feasibility"]) or {}
    prioritization = _load_latest_payload(workshop_id, ["results_prioritization"]) or {}
    action_plan = _load_latest_payload(workshop_id, ["results_action_plan"]) or {}
    discussion = _load_latest_payload(workshop_id, ["discussion"]) or {}

    # Ideas
    ideas = BrainstormIdea.query.options(joinedload(BrainstormIdea.participant)).filter(
        BrainstormIdea.task.has(workshop_id=workshop_id)
    ).order_by(BrainstormIdea.id.asc()).all()
    ideas_json = [
        {
            "idea_id": i.id,
            "participant_id": i.participant_id,
            "source": i.source,
            "text": i.corrected_text or i.content,
            "cluster_id": i.cluster_id,
            "duplicate_of_id": i.duplicate_of_id,
        }
        for i in ideas
    ]

    # Clusters with votes
    cl_rows = (
        db.session.query(
            IdeaCluster.id.label("cluster_id"),
            IdeaCluster.name,
            IdeaCluster.description,
            func.count(IdeaVote.id).label("votes"),
        )
        .outerjoin(IdeaVote, IdeaCluster.id == IdeaVote.cluster_id)
        .filter(IdeaCluster.task.has(workshop_id=workshop_id))
        .group_by(IdeaCluster.id)
        .order_by(func.count(IdeaVote.id).desc(), IdeaCluster.id.asc())
        .all()
    )
    clusters_json = []
    for r in cl_rows:
        clusters_json.append(
            {"cluster_id": int(r.cluster_id), "name": _safe(r.name), "summary": _safe(r.description), "votes": int(r.votes or 0)}
        )

    # Chat + Transcript snippets
    chats = ChatMessage.query.filter_by(workshop_id=workshop_id).order_by(ChatMessage.timestamp.asc()).all()
    chat_json = [{"user_id": c.user_id, "username": c.username, "message": c.message, "time": c.timestamp.isoformat()} for c in chats[-40:]]

    trans = Transcript.query.filter_by(workshop_id=workshop_id).order_by(Transcript.created_timestamp.asc()).all()
    transcripts_json = [
        {
            "transcript_id": t.transcript_id,
            "user_id": t.user_id,
            "entry_type": t.entry_type,
            "text": t.processed_transcript or t.raw_stt_transcript,
            "start": t.start_timestamp.isoformat() if t.start_timestamp else None,
            "end": t.end_timestamp.isoformat() if t.end_timestamp else None,
        }
        for t in trans[-80:]
    ]

    return {
        "workshop_overview": json.dumps(overview, ensure_ascii=False, indent=2),
        "pre_workshop_data": pre,
        "framing_json": json.dumps(framing, ensure_ascii=False, indent=2),
        "warmup_json": json.dumps(warmup, ensure_ascii=False, indent=2),
        "brainstorming_json": json.dumps(brainstorming, ensure_ascii=False, indent=2),
        "clustering_voting_json": json.dumps(clustering_voting, ensure_ascii=False, indent=2),
        "feasibility_json": json.dumps(feasibility, ensure_ascii=False, indent=2),
        "prioritization_json": json.dumps(prioritization, ensure_ascii=False, indent=2),
        "action_plan_json": json.dumps(action_plan, ensure_ascii=False, indent=2),
        "discussion_json": json.dumps(discussion, ensure_ascii=False, indent=2),
        "ideas_json": json.dumps(ideas_json, ensure_ascii=False, indent=2),
        "clusters_json": json.dumps(clusters_json, ensure_ascii=False, indent=2),
        "chat_json": json.dumps(chat_json, ensure_ascii=False, indent=2),
        "transcripts_json": json.dumps(transcripts_json, ensure_ascii=False, indent=2),
    }


def _parse_json_value(raw: Any) -> Any:
    if raw is None:
        return None
    if isinstance(raw, (dict, list)):
        return raw
    if isinstance(raw, str):
        text = raw.strip()
        if not text:
            return None
        try:
            return json.loads(text)
        except Exception:
            return None
    try:
        return json.loads(str(raw))
    except Exception:
        return None


def _build_canonical_session(inputs: Dict[str, Any]) -> Dict[str, Any]:
    canonical: Dict[str, Any] = {}

    def _assign(dest_key: str, source_key: str, default: Any) -> None:
        value = _parse_json_value(inputs.get(source_key))
        if value is None:
            value = default
        canonical[dest_key] = value

    _assign("workshop_overview", "workshop_overview", {})
    _assign("framing", "framing_json", {})
    _assign("warmup", "warmup_json", {})
    _assign("brainstorming", "brainstorming_json", {})
    _assign("clustering_voting", "clustering_voting_json", {})
    _assign("feasibility", "feasibility_json", {})
    _assign("prioritization", "prioritization_json", {})
    _assign("action_plan", "action_plan_json", {})
    _assign("discussion", "discussion_json", {})
    _assign("ideas", "ideas_json", [])
    _assign("clusters", "clusters_json", [])
    _assign("chat", "chat_json", [])
    _assign("transcripts", "transcripts_json", [])

    clusters_value = canonical.get("clusters")
    if not isinstance(clusters_value, list):
        clusters_list: List[Dict[str, Any]] = []
    else:
        clusters_list = [c for c in clusters_value if isinstance(c, dict)]
    canonical["clusters"] = clusters_list
    canonical["votes"] = [
        {"cluster_id": c.get("cluster_id"), "votes": c.get("votes")}
        for c in clusters_list
        if c.get("cluster_id") is not None
    ]

    canonical["pre_workshop_data"] = inputs.get("pre_workshop_data") or ""

    return canonical


# =============== LLM Invocation ===============

def _invoke_summary_model(inputs: Dict[str, Any]) -> Dict[str, Any]:
    llm = get_chat_llm_pro(model_kwargs={"temperature": 0.45, "max_tokens": 4000})
    template = """
You are the workshop composer and closing facilitator. Using ONLY the provided data, create a share-ready executive package for the group.

Return ONE strict JSON object with exactly these top-level keys (no markdown fences, no extra text). Every key below is required—do not omit any:

- title: "Workshop Summary"
- task_type: "summary"
- task_description: brief sentence describing the wrap-up step
- instructions: one paragraph guiding the group on how to use the artifacts
- task_duration: integer seconds for this phase
- narration: one paragraph, facilitator voice, friendly, inclusive
- tts_script: one paragraph (90–180 words), no lists, plain characters
- tts_read_time_seconds: integer ≥45
- artifacts: object with exactly these keys
    - markdown_doc: a concise one-page executive brief in Markdown (string). Sections: Executive Summary, Highlights, Decisions, Next Steps, Risks & Watch-outs.
    - document_spec: JSON for a professional PDF (cover+sections+appendices). Use sections that mirror a full report: Executive Summary, Session Highlights, Shortlist & Scores, Feasibility Findings, Decisions & Action Items, Risks & Mitigations, Timeline & Milestones, Appendix (clusters/ideas tables).
    - slides_spec: JSON for a slide deck (title + slides with layouts title+bullets/title+table). Include a cover, highlights, shortlist, feasibility, decisions, action plan, next steps.
- canonical_session_json: set this field to the literal string "AUTO" (the platform will insert normalized session data—do not embed large JSON here).

Hard rules:
- Use ONLY the data below; when unknown, write "TBD".
- Do NOT invent clusters or numbers; rely on clusters_json and prior payloads.
- JSON MUST be valid: no trailing commas; strings for Markdown; arrays for tables.
- Keep labels concise; speak plainly in narration/tts.

Workshop Snapshot:
{workshop_overview}

Framing:
{framing_json}

Warm-Up:
{warmup_json}

Brainstorming:
{brainstorming_json}

Clustering & Voting:
{clustering_voting_json}

Feasibility:
{feasibility_json}

Prioritization:
{prioritization_json}

Action Plan:
{action_plan_json}

Discussion:
{discussion_json}

Ideas:
{ideas_json}

Clusters:
{clusters_json}

Chat:
{chat_json}

Transcripts:
{transcripts_json}

Pre-Workshop Research (may be truncated):
{pre_workshop_data}
"""
    prompt = PromptTemplate.from_template(template)
    chain = prompt | llm
    raw = chain.invoke(inputs)
    text = _coerce_text(raw)
    block = extract_json_block(text) or text
    try:
        data = json.loads(block)
    except Exception as exc:
        raise SummaryGenerationError(f"Model did not return valid JSON: {exc}") from exc
    if not isinstance(data, dict):
        raise SummaryGenerationError("Model output must be a JSON object.")
    data["canonical_session_json"] = _build_canonical_session(inputs)
    return data


# =============== API entrypoint ===============

def get_summary_payload(workshop_id: int, phase_context: str) -> Dict[str, Any] | Tuple[str, int]:
    ws = db.session.get(Workshop, workshop_id)
    if not ws:
        return "Workshop not found", 404

    try:
        inputs = _collect_full_context(workshop_id)
    except SummaryGenerationError as exc:
        return str(exc), 400
    except Exception as exc:
        current_app.logger.error("[Summary] Input collection failed: %s", exc, exc_info=True)
        return "Failed to collect summary inputs", 500

    try:
        data = _invoke_summary_model(inputs)
    except SummaryGenerationError as exc:
        current_app.logger.error("[Summary] LLM failure: %s", exc, exc_info=True)
        return str(exc), 503
    except Exception as exc:
        current_app.logger.error("[Summary] Unhandled LLM error: %s", exc, exc_info=True)
        return "Summary generation error", 503

    artifacts = data.get("artifacts") if isinstance(data, dict) else None
    if isinstance(artifacts, dict) and "summary_report" not in data:
        markdown_doc = artifacts.get("markdown_doc")
        if isinstance(markdown_doc, str) and markdown_doc.strip():
            data["summary_report"] = markdown_doc

    # Strict: we persist EXACTLY what LLM returned (no rewrites/fallbacks)
    required = [
        "title", "task_type", "task_description", "instructions", "task_duration",
        "narration", "tts_script", "tts_read_time_seconds", "artifacts", "canonical_session_json",
        "summary_report"
    ]
    if not all(k in data for k in required):
        return "Summary output missing required fields", 500

    # Persist BrainstormTask with raw LLM payload
    task = BrainstormTask()
    task.workshop_id = workshop_id
    task.task_type = _safe(data.get("task_type") or "summary")
    task.title = _safe(data.get("title") or "Workshop Summary")
    task.description = _safe(data.get("task_description"))
    task.duration = int(data.get("task_duration") or 300)
    task.status = "pending"
    payload_str = json.dumps(data, ensure_ascii=False)
    task.prompt = payload_str
    task.payload_json = payload_str
    db.session.add(task)
    db.session.flush()
    data["task_id"] = task.id

    # Build artifacts from LLM-owned specs (if provided)
    try:
        art = data.get("artifacts") or {}
        # PDF
        if art.get("document_spec"):
            pdf_assets = _generate_pdf_from_doc_spec(ws, art["document_spec"])
            if pdf_assets:
                abs_path, rel_path, url = pdf_assets
                data["summary_pdf_path"] = rel_path
                data["summary_pdf_url"] = url
                _attach_doc(ws, abs_path, rel_path, url, data, "Summary Report")
        # PPTX
        if art.get("slides_spec"):
            pptx_assets = _generate_pptx_from_slides_spec(ws, art["slides_spec"])
            if pptx_assets:
                abs_path, rel_path, url = pptx_assets
                data["summary_pptx_path"] = rel_path
                data["summary_pptx_url"] = url
                _attach_doc(ws, abs_path, rel_path, url, data, "Summary Slides")
        # Markdown (store as a Document for convenience)
        if art.get("markdown_doc"):
            ts = datetime.utcnow().strftime("%Y-%m-%d %H-%M-%S")
            fname = f"{_safe_title(ws.title)} executive-brief {ts}.md"
            abs_path = os.path.join(_reports_dir(), fname)
            rel_path = os.path.join("uploads", "reports", fname)
            with open(abs_path, "w", encoding="utf-8") as f:
                f.write(str(art["markdown_doc"]))
            url = _media_url(rel_path)
            data["summary_markdown_path"] = rel_path
            data["summary_markdown_url"] = url
            _attach_doc(ws, abs_path, rel_path, url, data, "Executive Brief (Markdown)")

        # Mirror augmented payload (with artifact links) back into task
        augmented = json.dumps(data, ensure_ascii=False)
        task.payload_json = augmented
        task.prompt = augmented
    except Exception as exc:
        # If artifact generation fails, still return the LLM payload untouched (no content fallback)
        current_app.logger.error("[Summary] Artifact build error: %s", exc, exc_info=True)

    current_app.logger.info("[Summary] Created task %s for workshop %s", task.id, ws.id)
    return data
